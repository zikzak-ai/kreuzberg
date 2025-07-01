"""This module provides functions to extract textual content from files.

It includes vendored code:

- The extract PPTX logic is based on code vendored from `markitdown` to extract text from PPTX files.
    See: https://github.com/microsoft/markitdown/blob/main/src/markitdown/_markitdown.py
    Refer to the markitdown repository for it's license (MIT).
"""

from __future__ import annotations

import re
from contextlib import suppress
from html import escape
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING, ClassVar

import pptx
from anyio import Path as AsyncPath
from pptx.enum.shapes import MSO_SHAPE_TYPE

from kreuzberg._extractors._base import Extractor
from kreuzberg._mime_types import MARKDOWN_MIME_TYPE, POWER_POINT_MIME_TYPE
from kreuzberg._types import ExtractionResult
from kreuzberg._utils._string import normalize_spaces

if TYPE_CHECKING:  # pragma: no cover
    from pptx.presentation import Presentation

    from kreuzberg._types import Metadata


class PresentationExtractor(Extractor):
    """Extractor for PowerPoint (.pptx) files.

    This extractor processes PowerPoint presentations and converts their content into Markdown format.
    It handles slides, shapes, images, tables, and slide notes, preserving the structure and content
    of the presentation in a readable text format.

    The extractor provides both synchronous and asynchronous methods for processing files either
    from disk or from bytes in memory.
    """

    SUPPORTED_MIME_TYPES: ClassVar[set[str]] = {POWER_POINT_MIME_TYPE}

    async def extract_bytes_async(self, content: bytes) -> ExtractionResult:
        """Asynchronously extract content from PowerPoint file bytes.

        Args:
            content: Raw bytes of the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        return self._extract_pptx(content)

    async def extract_path_async(self, path: Path) -> ExtractionResult:
        """Asynchronously extract content from a PowerPoint file on disk.

        Args:
            path: Path to the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        content = await AsyncPath(path).read_bytes()
        return self._extract_pptx(content)

    def extract_bytes_sync(self, content: bytes) -> ExtractionResult:
        """Synchronously extract content from PowerPoint file bytes.

        Args:
            content: Raw bytes of the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        return self._extract_pptx(content)

    def extract_path_sync(self, path: Path) -> ExtractionResult:
        """Synchronously extract content from a PowerPoint file on disk.

        Args:
            path: Path to the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        content = Path(path).read_bytes()
        return self._extract_pptx(content)

    def _extract_pptx(self, file_contents: bytes) -> ExtractionResult:
        """Process PowerPoint file contents and convert to Markdown.

        This method handles the core logic of extracting content from a PowerPoint file.
        It processes:
        - Slide titles and content
        - Images (with alt text if available)
        - Tables (converted to HTML format)
        - Text frames
        - Slide notes

        Args:
            file_contents: Raw bytes of the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.

        Notes:
            The extraction preserves the following elements:
            - Slide numbers (as HTML comments)
            - Images (converted to Markdown image syntax with alt text)
            - Tables (converted to HTML table syntax)
            - Text content (with titles properly formatted)
            - Slide notes (under a dedicated section for each slide)
        """
        md_content = ""
        presentation = pptx.Presentation(BytesIO(file_contents))

        for index, slide in enumerate(presentation.slides):
            md_content += f"\n\n<!-- Slide number: {index + 1} -->\n"

            title = None
            if hasattr(slide.shapes, "title"):
                title = slide.shapes.title

            for shape in slide.shapes:
                if not hasattr(shape, "shape_type"):
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                    shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
                ):
                    alt_text = ""
                    with suppress(AttributeError):
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  # noqa: SLF001

                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    html_table = "<table>"
                    first_row = True

                    for row in shape.table.rows:
                        html_table += "<tr>"

                        for cell in row.cells:
                            tag = "th" if first_row else "td"
                            html_table += f"<{tag}>{escape(cell.text)}</{tag}>"

                        html_table += "</tr>"
                        first_row = False

                    html_table += "</table>"
                    md_content += "\n" + html_table + "\n"

                elif shape.has_text_frame:
                    md_content += "# " + shape.text.lstrip() + "\n" if shape == title else shape.text + "\n"

            md_content = md_content.strip()
            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame

                if notes_frame is not None:  # pragma: no branch
                    md_content += notes_frame.text

                md_content = md_content.strip()

        return ExtractionResult(
            content=normalize_spaces(md_content),
            mime_type=MARKDOWN_MIME_TYPE,
            metadata=self._extract_presentation_metadata(presentation),
            chunks=[],
        )

    @staticmethod
    def _extract_presentation_metadata(presentation: Presentation) -> Metadata:
        """Extract metadata from a presentation instance.

        Args:
            presentation: A `Presentation` object representing the PowerPoint file.

        Returns:
            PresentationMetadata: Object containing presentation-specific metadata fields.
        """
        metadata: Metadata = {}

        for metadata_key, core_property_key in [
            ("authors", "author"),
            ("comments", "comments"),
            ("status", "content_status"),
            ("created_by", "created"),
            ("identifier", "identifier"),
            ("keywords", "keywords"),
            ("modified_by", "last_modified_by"),
            ("modified_at", "modified"),
            ("version", "revision"),  # if version and revision are given, version overwrites
            ("subject", "subject"),
            ("title", "title"),
            ("version", "version"),
        ]:
            if core_property := getattr(presentation.core_properties, core_property_key, None):
                metadata[metadata_key] = core_property  # type: ignore[literal-required]

        if presentation.core_properties.language:
            metadata["languages"] = [presentation.core_properties.language]

        if presentation.core_properties.category:
            metadata["categories"] = [presentation.core_properties.category]

        fonts = set()
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, "font") and run.font.name:
                            fonts.add(run.font.name)

        if fonts:
            metadata["fonts"] = list(fonts)

        return metadata
