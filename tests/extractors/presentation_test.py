from __future__ import annotations

from typing import TYPE_CHECKING, Any

import pytest

from kreuzberg._extractors._presentation import PresentationExtractor
from kreuzberg.extraction import DEFAULT_CONFIG

if TYPE_CHECKING:
    from pathlib import Path

    from pytest_mock import MockerFixture


@pytest.fixture
def extractor() -> PresentationExtractor:
    return PresentationExtractor(
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", config=DEFAULT_CONFIG
    )


@pytest.mark.anyio
async def test_extract_pptx_with_notes(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_notes_slide = mocker.MagicMock()
    mock_text_frame = mocker.MagicMock()

    mock_presentation.slides = [mock_slide]
    mock_slide.has_notes_slide = True
    mock_slide.notes_slide = mock_notes_slide
    mock_notes_slide.notes_text_frame = mock_text_frame
    mock_text_frame.text = "Test note content"
    mock_slide.shapes = []

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = await extractor.extract_bytes_async(b"mock pptx content")

    assert "Test note content" in result.content
    assert result.mime_type == "text/markdown"


@pytest.mark.anyio
async def test_extract_path_async(mocker: MockerFixture, extractor: PresentationExtractor, tmp_path: Path) -> None:
    mock_presentation = mocker.MagicMock()
    mock_presentation.slides = []
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    test_file = tmp_path / "test.pptx"
    test_file.write_bytes(b"mock pptx content")

    result = await extractor.extract_path_async(test_file)

    assert result.mime_type == "text/markdown"


def test_extract_bytes_sync(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_presentation.slides = []
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert result.mime_type == "text/markdown"


def test_extract_path_sync(mocker: MockerFixture, extractor: PresentationExtractor, tmp_path: Path) -> None:
    mock_presentation = mocker.MagicMock()
    mock_presentation.slides = []
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    test_file = tmp_path / "test.pptx"
    test_file.write_bytes(b"mock pptx content")

    result = extractor.extract_path_sync(test_file)

    assert result.mime_type == "text/markdown"


def test_extract_pptx_with_slide_title(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_title_shape = mocker.MagicMock()
    mock_shapes = mocker.MagicMock()

    mock_shapes.title = mock_title_shape
    mock_slide.shapes = mock_shapes
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mock_shapes.__iter__ = mocker.Mock(return_value=iter([mock_title_shape]))
    mock_title_shape.shape_type = mocker.MagicMock()
    mock_title_shape.has_text_frame = True
    mock_title_shape.text = "Test Title"

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "# Test Title" in result.content
    assert result.mime_type == "text/markdown"


def test_extract_pptx_with_shapes_no_shape_type(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_shape_no_type = mocker.MagicMock()

    del mock_shape_no_type.shape_type

    mock_slide.shapes = [mock_shape_no_type]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert result.mime_type == "text/markdown"


def test_extract_pptx_with_picture_shape(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_picture_shape = mocker.MagicMock()

    mock_picture_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    mock_picture_shape.name = "test_image"

    mock_element = mocker.MagicMock()
    mock_nvXxPr = mocker.MagicMock()
    mock_cNvPr = mocker.MagicMock()
    mock_cNvPr.attrib = {"descr": "Test alt text"}
    mock_nvXxPr.cNvPr = mock_cNvPr
    mock_element._nvXxPr = mock_nvXxPr
    mock_picture_shape._element = mock_element

    mock_slide.shapes = [mock_picture_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "![Test alt text](test_image.jpg)" in result.content
    assert result.mime_type == "text/markdown"


def test_extract_pptx_with_table_shape(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_table_shape = mocker.MagicMock()

    mock_table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

    mock_table = mocker.MagicMock()
    mock_row1 = mocker.MagicMock()
    mock_row2 = mocker.MagicMock()

    mock_cell1 = mocker.MagicMock()
    mock_cell1.text = "Header1"
    mock_cell2 = mocker.MagicMock()
    mock_cell2.text = "Header2"
    mock_cell3 = mocker.MagicMock()
    mock_cell3.text = "Data1"
    mock_cell4 = mocker.MagicMock()
    mock_cell4.text = "Data2"

    mock_row1.cells = [mock_cell1, mock_cell2]
    mock_row2.cells = [mock_cell3, mock_cell4]
    mock_table.rows = [mock_row1, mock_row2]
    mock_table_shape.table = mock_table

    mock_slide.shapes = [mock_table_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "<table>" in result.content
    assert "<th>Header1</th>" in result.content
    assert "<td>Data1</td>" in result.content
    assert result.mime_type == "text/markdown"


def test_extract_pptx_with_text_frame_shape(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_text_shape = mocker.MagicMock()

    mock_text_shape.shape_type = mocker.MagicMock()
    mock_text_shape.has_text_frame = True
    mock_text_shape.text = "  Some text content"

    mock_slide.shapes = [mock_text_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "Some text content" in result.content
    assert result.mime_type == "text/markdown"


def test_extract_presentation_metadata(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = "Test Author"
    mock_core_properties.title = "Test Title"
    mock_core_properties.subject = "Test Subject"
    mock_core_properties.language = "en-US"
    mock_core_properties.category = "Test Category"
    mock_core_properties.comments = None
    mock_core_properties.content_status = None
    mock_core_properties.created = None
    mock_core_properties.identifier = None
    mock_core_properties.keywords = None
    mock_core_properties.last_modified_by = None
    mock_core_properties.modified = None
    mock_core_properties.revision = None
    mock_core_properties.version = None

    mock_presentation.core_properties = mock_core_properties

    mock_shape = mocker.MagicMock()
    mock_text_frame = mocker.MagicMock()
    mock_paragraph = mocker.MagicMock()
    mock_run = mocker.MagicMock()
    mock_font = mocker.MagicMock()
    mock_font.name = "Arial"
    mock_run.font = mock_font
    mock_paragraph.runs = [mock_run]
    mock_text_frame.paragraphs = [mock_paragraph]
    mock_shape.text_frame = mock_text_frame

    mock_slide.shapes = [mock_shape]
    mock_presentation.slides = [mock_slide]

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert result.metadata["authors"] == "Test Author"  # type: ignore[comparison-overlap]
    assert result.metadata["title"] == "Test Title"
    assert result.metadata["subject"] == "Test Subject"
    assert result.metadata["languages"] == ["en-US"]
    assert result.metadata["categories"] == ["Test Category"]
    assert result.metadata["fonts"] == ["Arial"]
    assert result.mime_type == "text/markdown"


def test_extract_presentation_metadata_no_text_frame(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = None
    mock_core_properties.title = None
    mock_core_properties.language = None
    mock_core_properties.category = None
    mock_presentation.core_properties = mock_core_properties

    mock_shape = mocker.MagicMock()
    del mock_shape.text_frame

    mock_slide.shapes = [mock_shape]
    mock_presentation.slides = [mock_slide]

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "fonts" not in result.metadata
    assert result.mime_type == "text/markdown"


def test_extract_pptx_shape_without_text_frame(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_shape = mocker.MagicMock()

    mock_shape.shape_type = mocker.MagicMock()
    mock_shape.has_text_frame = False

    mock_slide.shapes = [mock_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert result.mime_type == "text/markdown"


def test_extract_presentation_metadata_run_without_font(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = None
    mock_core_properties.language = None
    mock_core_properties.category = None
    mock_presentation.core_properties = mock_core_properties

    mock_shape = mocker.MagicMock()
    mock_text_frame = mocker.MagicMock()
    mock_paragraph = mocker.MagicMock()
    mock_run_no_font = mocker.MagicMock()

    del mock_run_no_font.font

    mock_paragraph.runs = [mock_run_no_font]
    mock_text_frame.paragraphs = [mock_paragraph]
    mock_shape.text_frame = mock_text_frame

    mock_slide.shapes = [mock_shape]
    mock_presentation.slides = [mock_slide]

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "fonts" not in result.metadata
    assert result.mime_type == "text/markdown"


def test_extract_presentation_metadata_font_without_name(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = None
    mock_core_properties.language = None
    mock_core_properties.category = None
    mock_presentation.core_properties = mock_core_properties

    mock_shape = mocker.MagicMock()
    mock_text_frame = mocker.MagicMock()
    mock_paragraph = mocker.MagicMock()
    mock_run = mocker.MagicMock()
    mock_font = mocker.MagicMock()
    mock_font.name = None
    mock_run.font = mock_font
    mock_paragraph.runs = [mock_run]
    mock_text_frame.paragraphs = [mock_paragraph]
    mock_shape.text_frame = mock_text_frame

    mock_slide.shapes = [mock_shape]
    mock_presentation.slides = [mock_slide]

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "fonts" not in result.metadata
    assert result.mime_type == "text/markdown"


def test_presentation_comprehensive_picture_shape_no_alt_text(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_picture_shape = mocker.MagicMock()

    mock_picture_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    mock_picture_shape.name = "test_image_no_alt"

    mock_element = mocker.MagicMock()
    del mock_element._nvXxPr
    mock_picture_shape._element = mock_element

    mock_slide.shapes = [mock_picture_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "![test_image_no_alt](test_image_no_alt.jpg)" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_comprehensive_placeholder_with_image(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_placeholder_shape = mocker.MagicMock()

    mock_placeholder_shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
    mock_placeholder_shape.name = "placeholder_image"
    mock_placeholder_shape.image = mocker.MagicMock()

    mock_element = mocker.MagicMock()
    mock_nvXxPr = mocker.MagicMock()
    mock_cNvPr = mocker.MagicMock()
    mock_cNvPr.attrib = {"descr": "Placeholder alt text"}
    mock_nvXxPr.cNvPr = mock_cNvPr
    mock_element._nvXxPr = mock_nvXxPr
    mock_placeholder_shape._element = mock_element

    mock_slide.shapes = [mock_placeholder_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "![Placeholder alt text](placeholder_image.jpg)" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_comprehensive_non_title_text_frame(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_title_shape = mocker.MagicMock()
    mock_text_shape = mocker.MagicMock()
    mock_shapes = mocker.MagicMock()

    mock_shapes.title = mock_title_shape
    mock_slide.shapes = mock_shapes
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mock_title_shape.shape_type = mocker.MagicMock()
    mock_title_shape.has_text_frame = True
    mock_title_shape.text = "  Title Text"

    mock_text_shape.shape_type = mocker.MagicMock()
    mock_text_shape.has_text_frame = True
    mock_text_shape.text = "Regular content text"

    mock_shapes.__iter__ = mocker.Mock(return_value=iter([mock_title_shape, mock_text_shape]))

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "# Title Text" in result.content
    assert "Regular content text" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_comprehensive_notes_text_frame_none(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_notes_slide = mocker.MagicMock()

    mock_presentation.slides = [mock_slide]
    mock_slide.has_notes_slide = True
    mock_slide.notes_slide = mock_notes_slide
    mock_notes_slide.notes_text_frame = None
    mock_slide.shapes = []

    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "### Notes:" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_comprehensive_shapes_no_title_attribute(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_text_shape = mocker.MagicMock()

    mock_shapes = mocker.MagicMock()
    del mock_shapes.title
    mock_slide.shapes = mock_shapes
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]

    mock_text_shape.shape_type = mocker.MagicMock()
    mock_text_shape.has_text_frame = True
    mock_text_shape.text = "Content without title"

    mock_shapes.__iter__ = mocker.Mock(return_value=iter([mock_text_shape]))

    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "Content without title" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_comprehensive_multiple_slides(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()

    mock_slide1 = mocker.MagicMock()
    mock_slide2 = mocker.MagicMock()
    mock_slide3 = mocker.MagicMock()

    mock_presentation.slides = [mock_slide1, mock_slide2, mock_slide3]

    mock_slide1.shapes = []
    mock_slide1.has_notes_slide = False

    mock_text_shape = mocker.MagicMock()
    mock_text_shape.shape_type = mocker.MagicMock()
    mock_text_shape.has_text_frame = True
    mock_text_shape.text = "Slide 2 content"
    mock_slide2.shapes = [mock_text_shape]
    mock_slide2.has_notes_slide = False

    mock_slide3.shapes = []
    mock_slide3.has_notes_slide = True
    mock_notes_slide = mocker.MagicMock()
    mock_notes_text_frame = mocker.MagicMock()
    mock_notes_text_frame.text = "Notes for slide 3"
    mock_notes_slide.notes_text_frame = mock_notes_text_frame
    mock_slide3.notes_slide = mock_notes_slide

    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "<!-- Slide number: 1 -->" in result.content
    assert "<!-- Slide number: 2 -->" in result.content
    assert "<!-- Slide number: 3 -->" in result.content
    assert "Slide 2 content" in result.content
    assert "### Notes:" in result.content
    assert "Notes for slide 3" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_metadata_all_properties(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    from datetime import datetime, timezone

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = "Full Author"
    mock_core_properties.comments = "Full Comments"
    mock_core_properties.content_status = "Final"
    mock_core_properties.created = datetime(2023, 1, 15, 10, 30, 0, tzinfo=timezone.utc)
    mock_core_properties.identifier = "PRES-12345"
    mock_core_properties.keywords = "presentation, test, comprehensive"
    mock_core_properties.last_modified_by = "Editor User"
    mock_core_properties.modified = datetime(2023, 2, 20, 14, 45, 0, tzinfo=timezone.utc)
    mock_core_properties.revision = 5
    mock_core_properties.subject = "Comprehensive Test Subject"
    mock_core_properties.title = "Comprehensive Test Title"
    mock_core_properties.language = "en-GB"
    mock_core_properties.category = "Educational"

    mock_presentation.core_properties = mock_core_properties

    mock_shape = mocker.MagicMock()
    mock_text_frame = mocker.MagicMock()
    mock_paragraph = mocker.MagicMock()
    mock_run = mocker.MagicMock()
    mock_font = mocker.MagicMock()
    mock_font.name = "Calibri"
    mock_run.font = mock_font
    mock_paragraph.runs = [mock_run]
    mock_text_frame.paragraphs = [mock_paragraph]
    mock_shape.text_frame = mock_text_frame

    mock_slide.shapes = [mock_shape]
    mock_slide.has_notes_slide = True
    mock_presentation.slides = [mock_slide]

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    metadata = result.metadata
    assert str(metadata["authors"]) == "Full Author"
    assert str(metadata["comments"]) == "Full Comments"
    assert str(metadata["status"]) == "Final"
    assert metadata["created_by"] == datetime(2023, 1, 15, 10, 30, 0, tzinfo=timezone.utc)
    assert str(metadata["identifier"]) == "PRES-12345"
    assert str(metadata["keywords"]) == "presentation, test, comprehensive"
    assert str(metadata["modified_by"]) == "Editor User"
    assert metadata["modified_at"] == datetime(2023, 2, 20, 14, 45, 0, tzinfo=timezone.utc)
    assert int(metadata["version"]) == 5
    assert metadata["subject"] == "Comprehensive Test Subject"
    assert metadata["title"] == "Comprehensive Test Title"
    assert metadata["languages"] == ["en-GB"]
    assert metadata["categories"] == ["Educational"]
    assert metadata["fonts"] == ["Calibri"]

    assert metadata["description"] == "Presentation with 1 slide, 1 with notes"
    assert "PowerPoint presentation with 1 slide" in metadata["summary"]
    assert "1 slides have notes" in metadata["summary"]
    assert "uses 1 font" in metadata["summary"]


def test_presentation_metadata_empty_presentation(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_presentation.slides = []
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "description" not in result.metadata
    assert "summary" not in result.metadata
    assert result.mime_type == "text/markdown"


def test_presentation_metadata_multiple_fonts(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = None
    mock_core_properties.language = None
    mock_core_properties.category = None
    mock_presentation.core_properties = mock_core_properties

    mock_shape1 = mocker.MagicMock()
    mock_text_frame1 = mocker.MagicMock()
    mock_paragraph1 = mocker.MagicMock()
    mock_run1 = mocker.MagicMock()
    mock_font1 = mocker.MagicMock()
    mock_font1.name = "Arial"
    mock_run1.font = mock_font1
    mock_paragraph1.runs = [mock_run1]
    mock_text_frame1.paragraphs = [mock_paragraph1]
    mock_shape1.text_frame = mock_text_frame1

    mock_shape2 = mocker.MagicMock()
    mock_text_frame2 = mocker.MagicMock()
    mock_paragraph2 = mocker.MagicMock()
    mock_run2a = mocker.MagicMock()
    mock_font2a = mocker.MagicMock()
    mock_font2a.name = "Calibri"
    mock_run2a.font = mock_font2a
    mock_run2b = mocker.MagicMock()
    mock_font2b = mocker.MagicMock()
    mock_font2b.name = "Arial"
    mock_run2b.font = mock_font2b
    mock_paragraph2.runs = [mock_run2a, mock_run2b]
    mock_text_frame2.paragraphs = [mock_paragraph2]
    mock_shape2.text_frame = mock_text_frame2

    mock_slide.shapes = [mock_shape1, mock_shape2]
    mock_presentation.slides = [mock_slide]

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert set(result.metadata["fonts"]) == {"Arial", "Calibri"}
    assert len(result.metadata["fonts"]) == 2
    assert "uses 2 fonts" in result.metadata["summary"]


def test_presentation_metadata_structure_single_vs_plural(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()

    mock_slide1 = mocker.MagicMock()
    mock_slide1.shapes = []
    mock_slide1.has_notes_slide = True

    mock_slide2 = mocker.MagicMock()
    mock_slide2.shapes = []
    mock_slide2.has_notes_slide = False

    mock_slide3 = mocker.MagicMock()
    mock_slide3.shapes = []
    mock_slide3.has_notes_slide = True

    mock_presentation.slides = [mock_slide1, mock_slide2, mock_slide3]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert result.metadata["description"] == "Presentation with 3 slides, 2 with notes"
    assert "PowerPoint presentation with 3 slides" in result.metadata["summary"]
    assert "2 slides have notes" in result.metadata["summary"]


def test_presentation_metadata_existing_summary_preserved(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_slide.shapes = []
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]

    mock_core_properties = mocker.MagicMock()
    mock_core_properties.author = None
    mock_core_properties.comments = "Existing summary content"
    mock_presentation.core_properties = mock_core_properties

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    original_extract_method = extractor._extract_presentation_metadata

    def mock_extract_with_existing_summary(presentation: Any) -> Any:
        metadata = original_extract_method(presentation)
        metadata["summary"] = "Pre-existing summary"
        return metadata

    mocker.patch.object(extractor, "_extract_presentation_metadata", side_effect=mock_extract_with_existing_summary)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert result.metadata["summary"] == "Pre-existing summary"


def test_presentation_table_processing_html_content(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_table_shape = mocker.MagicMock()

    mock_table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

    mock_table = mocker.MagicMock()
    mock_row1 = mocker.MagicMock()
    mock_row2 = mocker.MagicMock()

    mock_cell1 = mocker.MagicMock()
    mock_cell1.text = "Header <with> special & chars"
    mock_cell2 = mocker.MagicMock()
    mock_cell2.text = 'Header "with" quotes'
    mock_cell3 = mocker.MagicMock()
    mock_cell3.text = "Data <script>alert('xss')</script>"
    mock_cell4 = mocker.MagicMock()
    mock_cell4.text = 'Data & more "special" chars'

    mock_row1.cells = [mock_cell1, mock_cell2]
    mock_row2.cells = [mock_cell3, mock_cell4]
    mock_table.rows = [mock_row1, mock_row2]
    mock_table_shape.table = mock_table

    mock_slide.shapes = [mock_table_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "&lt;with&gt;" in result.content
    assert "&amp;" in result.content
    assert "&quot;with&quot;" in result.content
    assert "&lt;script&gt;" in result.content
    assert "<table>" in result.content
    assert "<th>" in result.content
    assert "<td>" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_table_processing_empty_cells(mocker: MockerFixture, extractor: PresentationExtractor) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_table_shape = mocker.MagicMock()

    mock_table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

    mock_table = mocker.MagicMock()
    mock_row1 = mocker.MagicMock()
    mock_row2 = mocker.MagicMock()

    mock_cell1 = mocker.MagicMock()
    mock_cell1.text = ""
    mock_cell2 = mocker.MagicMock()
    mock_cell2.text = "Header2"
    mock_cell3 = mocker.MagicMock()
    mock_cell3.text = "Data1"
    mock_cell4 = mocker.MagicMock()
    mock_cell4.text = ""

    mock_row1.cells = [mock_cell1, mock_cell2]
    mock_row2.cells = [mock_cell3, mock_cell4]
    mock_table.rows = [mock_row1, mock_row2]
    mock_table_shape.table = mock_table

    mock_slide.shapes = [mock_table_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "<th></th>" in result.content
    assert "<th>Header2</th>" in result.content
    assert "<td>Data1</td>" in result.content
    assert "<td></td>" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_image_processing_special_characters_in_name(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_picture_shape = mocker.MagicMock()

    mock_picture_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    mock_picture_shape.name = "test-image_with@special#chars!"

    mock_element = mocker.MagicMock()
    mock_nvXxPr = mocker.MagicMock()
    mock_cNvPr = mocker.MagicMock()
    mock_cNvPr.attrib = {"descr": "Image with special chars"}
    mock_nvXxPr.cNvPr = mock_cNvPr
    mock_element._nvXxPr = mock_nvXxPr
    mock_picture_shape._element = mock_element

    mock_slide.shapes = [mock_picture_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "![Image with special chars](testimage_withspecialchars.jpg)" in result.content
    assert result.mime_type == "text/markdown"


def test_presentation_image_processing_placeholder_without_image_attribute(
    mocker: MockerFixture, extractor: PresentationExtractor
) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_presentation = mocker.MagicMock()
    mock_slide = mocker.MagicMock()
    mock_placeholder_shape = mocker.MagicMock()

    mock_placeholder_shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
    mock_placeholder_shape.name = "text_placeholder"
    mock_placeholder_shape.has_text_frame = True
    mock_placeholder_shape.text = "Placeholder text content"
    del mock_placeholder_shape.image

    mock_slide.shapes = [mock_placeholder_shape]
    mock_slide.has_notes_slide = False
    mock_presentation.slides = [mock_slide]
    mock_presentation.core_properties = mocker.MagicMock()
    mock_presentation.core_properties.author = None

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = extractor.extract_bytes_sync(b"mock pptx content")

    assert "Placeholder text content" in result.content
    assert "![" not in result.content
    assert result.mime_type == "text/markdown"
