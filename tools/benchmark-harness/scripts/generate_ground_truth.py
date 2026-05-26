#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "beautifulsoup4>=4.12",
#     "python-docx>=1.0",
#     "python-pptx>=1.0",
#     "openpyxl>=3.1",
#     "nbformat>=5.9",
#     "xlrd>=2.0",
#     "extract-msg>=0.48",
#     "lxml>=5.0",
#     "odfpy>=1.4",
# ]
# ///
"""Generate ground truth text files for benchmark fixtures.

Walks all fixture JSONs, extracts text from source documents using independent
tools (not benchmarked frameworks), writes ground truth .txt files, patches
fixture JSONs with ground_truth field, and updates ground_truth_mapping.json.

PDF Ground Truth Methodology (updated Feb 2026):
    PDF ground truth was regenerated using AI visual extraction (Claude Haiku
    reading each PDF page as an image) for scanned/complex PDFs, and pdftotext
    for born-digital PDFs with reliable embedded text. The previous approach of
    using pdftotext for all PDFs produced incorrect ground truth for scanned
    documents since pdftotext cannot read image-based text.

    The handle_pdftotext() function below is retained for regenerating GT from
    born-digital PDFs. For scanned PDFs, GT files were manually curated via AI
    extraction and should not be overwritten by running this script with --force.

Usage:
    uv run tools/benchmark-harness/scripts/generate_ground_truth.py [OPTIONS]

Options:
    --dry-run           Print planned actions without writing
    --format-filter     Comma-separated file types to process (e.g., md,txt,pdf)
    --force             Regenerate even if ground truth already exists
    --skip-types        Comma-separated file types to skip
"""

from __future__ import annotations

import argparse
import email
import json
import os
import subprocess
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

# ---------------------------------------------------------------------------
# File type → handler mapping
# ---------------------------------------------------------------------------

RAW_SOURCE_TYPES = frozenset(
    {
        "md",
        "txt",
        "rst",
        "org",
        "commonmark",
        "djot",
        "toml",
        "yaml",
        "json",
        "tsv",
        "bib",
        "csv",
        "svg",
    }
)

PDFTOTEXT_TYPES = frozenset({"pdf"})
PANDOC_TYPES = frozenset(
    {
        "tex",
        "latex",
        "typ",
        "epub",
        "fb2",
        "docbook",
        "odt",
        "rtf",
        "opml",
    }
)
PYTHON_DOCX_TYPES = frozenset({"docx"})
PYTHON_PPTX_TYPES = frozenset({"pptx", "pptm", "ppsx"})
OPENPYXL_TYPES = frozenset({"xlsx", "xlsm"})
ODS_TYPES = frozenset({"ods"})
BEAUTIFULSOUP_TYPES = frozenset({"html"})
PYTHON_EMAIL_TYPES = frozenset({"eml"})
EXTRACT_MSG_TYPES = frozenset({"msg"})
NBFORMAT_TYPES = frozenset({"ipynb"})
XML_PARSE_TYPES = frozenset({"xml"})
XLRD_TYPES = frozenset({"xls"})
ANTIWORD_TYPES = frozenset({"doc"})
LIBREOFFICE_TYPES = frozenset({"ppt"})
DBF_TYPES = frozenset({"dbf"})
HWP_TYPES = frozenset({"hwp"})

# Archive and image types are excluded from ground truth generation
EXCLUDED_TYPES = frozenset(
    {
        "7z",
        "gz",
        "tar",
        "tgz",
        "zip",
        "lz4",
        "gif",
        "jpeg",
        "jpg",
        "jp2",
        "png",
        "tiff",
        "webp",
        "bmp",
        "pbm",
        "pgm",
        "pnm",
        "ppm",
    }
)

ALL_HANDLED_TYPES = (
    RAW_SOURCE_TYPES
    | PDFTOTEXT_TYPES
    | PANDOC_TYPES
    | PYTHON_DOCX_TYPES
    | PYTHON_PPTX_TYPES
    | OPENPYXL_TYPES
    | BEAUTIFULSOUP_TYPES
    | PYTHON_EMAIL_TYPES
    | EXTRACT_MSG_TYPES
    | NBFORMAT_TYPES
    | XML_PARSE_TYPES
    | XLRD_TYPES
    | ANTIWORD_TYPES
    | LIBREOFFICE_TYPES
    | ODS_TYPES
    | DBF_TYPES
    | HWP_TYPES
)


def get_source_type(file_type: str) -> str:
    """Return the ground truth source type string for a given file type."""
    if file_type in RAW_SOURCE_TYPES:
        return "raw_source"
    if file_type in PDFTOTEXT_TYPES:
        return "pdftotext"
    if file_type in PANDOC_TYPES:
        return "pandoc"
    if file_type in PYTHON_DOCX_TYPES:
        return "python-docx"
    if file_type in PYTHON_PPTX_TYPES:
        return "python-pptx"
    if file_type in OPENPYXL_TYPES:
        return "openpyxl"
    if file_type in BEAUTIFULSOUP_TYPES:
        return "beautifulsoup"
    if file_type in PYTHON_EMAIL_TYPES:
        return "python_email"
    if file_type in EXTRACT_MSG_TYPES:
        return "extract_msg"
    if file_type in NBFORMAT_TYPES:
        return "nbformat"
    if file_type in XML_PARSE_TYPES:
        return "xml_parse"
    if file_type in XLRD_TYPES:
        return "xlrd"
    if file_type in ANTIWORD_TYPES:
        return "antiword"
    if file_type in LIBREOFFICE_TYPES:
        return "libreoffice"
    if file_type in ODS_TYPES:
        return "odfpy"
    if file_type in DBF_TYPES:
        return "manual"
    if file_type in HWP_TYPES:
        return "manual"
    return "manual"


# ---------------------------------------------------------------------------
# Text extraction handlers
# ---------------------------------------------------------------------------


def handle_raw_source(doc_path: Path) -> str:
    """Read the file as-is. For text-based formats, source content IS ground truth."""
    try:
        return doc_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return doc_path.read_text(encoding="latin-1")


def handle_pdftotext(doc_path: Path) -> str:
    """Extract text from PDF using pdftotext (poppler-utils).

    Note: This works well for born-digital PDFs with embedded text layers.
    For scanned PDFs, pdftotext produces garbage output. Scanned PDF ground
    truth should be generated via AI visual extraction instead.
    """
    result = subprocess.run(
        ["pdftotext", "-layout", str(doc_path), "-"],
        capture_output=True,
        text=True,
        timeout=60,
    )
    if result.returncode != 0:
        raise RuntimeError(f"pdftotext failed: {result.stderr}")
    return result.stdout


def handle_pandoc(doc_path: Path, file_type: str) -> str:
    """Convert document to plain text using pandoc."""
    # Map file types to pandoc input formats
    pandoc_format_map = {
        "tex": "latex",
        "latex": "latex",
        "typ": "typst",
        "epub": "epub",
        "fb2": "fb2",
        "docbook": "docbook",
        "odt": "odt",
        "rtf": "rtf",
        "opml": "opml",
        "doc": "doc",
        "ppt": "ppt",
    }
    input_format = pandoc_format_map.get(file_type)
    cmd = ["pandoc", "-t", "plain", "--wrap=none", str(doc_path)]
    if input_format:
        cmd.insert(1, "-f")
        cmd.insert(2, input_format)
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f"pandoc failed: {result.stderr}")
    return result.stdout


def handle_python_docx(doc_path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    import docx

    doc = docx.Document(str(doc_path))
    paragraphs = [p.text for p in doc.paragraphs]
    # Also extract table text
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            paragraphs.append("\t".join(cells))
    return "\n".join(paragraphs)


def handle_python_pptx(doc_path: Path) -> str:
    """Extract text from PPTX/PPTM/PPSX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(doc_path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def handle_openpyxl(doc_path: Path) -> str:
    """Extract text from XLSX/XLSM using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(str(doc_path), read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def handle_beautifulsoup(doc_path: Path) -> str:
    """Extract text from HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup

    try:
        html_content = doc_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        html_content = doc_path.read_text(encoding="latin-1")
    soup = BeautifulSoup(html_content, "html.parser")
    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def handle_python_email(doc_path: Path) -> str:
    """Extract text from EML using Python email stdlib."""
    try:
        raw = doc_path.read_bytes()
        msg = email.message_from_bytes(raw)
    except Exception:
        raw = doc_path.read_text(encoding="utf-8", errors="replace")
        msg = email.message_from_string(raw)

    parts = []
    # Add headers
    for header in ("From", "To", "Subject", "Date"):
        val = msg.get(header)
        if val:
            parts.append(f"{header}: {val}")

    if parts:
        parts.append("")  # blank line after headers

    # Extract body
    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    try:
                        parts.append(payload.decode(charset, errors="replace"))
                    except (LookupError, UnicodeDecodeError):
                        parts.append(payload.decode("utf-8", errors="replace"))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            try:
                parts.append(payload.decode(charset, errors="replace"))
            except (LookupError, UnicodeDecodeError):
                parts.append(payload.decode("utf-8", errors="replace"))

    return "\n".join(parts)


def handle_extract_msg(doc_path: Path) -> str:
    """Extract text from MSG using extract-msg."""
    import extract_msg

    msg = extract_msg.openMsg(str(doc_path))
    parts = []
    if msg.subject:
        parts.append(f"Subject: {msg.subject}")
    if msg.sender:
        parts.append(f"From: {msg.sender}")
    if msg.to:
        parts.append(f"To: {msg.to}")
    if msg.date:
        parts.append(f"Date: {msg.date}")
    if parts:
        parts.append("")
    if msg.body:
        parts.append(msg.body)
    msg.close()
    return "\n".join(parts)


def handle_nbformat(doc_path: Path) -> str:
    """Extract text from Jupyter notebooks using nbformat."""
    import nbformat

    nb = nbformat.read(str(doc_path), as_version=4)
    parts = []
    for cell in nb.cells:
        if cell.cell_type in ("code", "markdown", "raw"):
            source = cell.source.strip()
            if source:
                parts.append(source)
    return "\n\n".join(parts)


def handle_xml_parse(doc_path: Path) -> str:
    """Extract text content from XML using xml.etree."""
    try:
        tree = ET.parse(str(doc_path))
    except ET.ParseError:
        # Fallback: read as raw text
        return handle_raw_source(doc_path)
    root = tree.getroot()
    texts = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            texts.append(elem.text.strip())
        if elem.tail and elem.tail.strip():
            texts.append(elem.tail.strip())
    return "\n".join(texts)


def handle_xlrd(doc_path: Path) -> str:
    """Extract text from XLS using xlrd."""
    import xlrd

    wb = xlrd.open_workbook(str(doc_path))
    lines = []
    for sheet_idx in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_idx)
        for row_idx in range(ws.nrows):
            cells = [str(ws.cell_value(row_idx, col_idx)) for col_idx in range(ws.ncols)]
            if any(c for c in cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


def handle_antiword(doc_path: Path) -> str:
    """Extract text from DOC using antiword, catdoc, or pandoc as fallbacks."""
    # Try antiword first
    try:
        result = subprocess.run(
            ["antiword", str(doc_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0:
            return result.stdout
    except FileNotFoundError:
        pass

    # Fallback to catdoc
    try:
        result = subprocess.run(
            ["catdoc", str(doc_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0:
            return result.stdout
    except FileNotFoundError:
        pass

    # Fallback to textutil (macOS)
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(doc_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0:
            return result.stdout
    except FileNotFoundError:
        pass

    raise RuntimeError("No DOC extraction tool available (need antiword, catdoc, or textutil)")


def handle_ods(doc_path: Path) -> str:
    """Extract text from ODS using odfpy."""
    from odf import text as odf_text
    from odf.opendocument import load as odf_load
    from odf.table import Table, TableCell, TableRow

    doc = odf_load(str(doc_path))
    lines = []
    for table in doc.spreadsheet.getElementsByType(Table):
        for row in table.getElementsByType(TableRow):
            cells = []
            for cell in row.getElementsByType(TableCell):
                # Get text content from cell
                cell_texts = []
                for p in cell.getElementsByType(odf_text.P):
                    # Recursively get all text
                    text_parts = []
                    for node in p.childNodes:
                        if hasattr(node, "data"):
                            text_parts.append(node.data)
                        elif hasattr(node, "__str__"):
                            text_parts.append(str(node))
                    cell_texts.append("".join(text_parts))
                # Handle repeated cells
                repeat = cell.getAttribute("numbercolumnsrepeated")
                cell_text = " ".join(cell_texts)
                if repeat and int(repeat) > 1 and cell_text:
                    cells.extend([cell_text] * min(int(repeat), 100))
                else:
                    cells.append(cell_text)
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


def handle_libreoffice(doc_path: Path) -> str:
    """Extract text from PPT using LibreOffice CLI, with pandoc fallback."""
    import tempfile

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, str(doc_path)],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode == 0:
                txt_files = list(Path(tmpdir).glob("*.txt"))
                if txt_files:
                    return txt_files[0].read_text(encoding="utf-8", errors="replace")
    except FileNotFoundError:
        pass

    # Fallback: try textutil (macOS)
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(doc_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0:
            return result.stdout
    except FileNotFoundError:
        pass

    raise RuntimeError("No PPT extraction tool available (need libreoffice or textutil)")


def extract_text(doc_path: Path, file_type: str) -> str:
    """Dispatch to the appropriate handler for the given file type."""
    if file_type in RAW_SOURCE_TYPES:
        return handle_raw_source(doc_path)
    if file_type in PDFTOTEXT_TYPES:
        return handle_pdftotext(doc_path)
    if file_type in PANDOC_TYPES:
        return handle_pandoc(doc_path, file_type)
    if file_type in PYTHON_DOCX_TYPES:
        return handle_python_docx(doc_path)
    if file_type in PYTHON_PPTX_TYPES:
        return handle_python_pptx(doc_path)
    if file_type in OPENPYXL_TYPES:
        return handle_openpyxl(doc_path)
    if file_type in BEAUTIFULSOUP_TYPES:
        return handle_beautifulsoup(doc_path)
    if file_type in PYTHON_EMAIL_TYPES:
        return handle_python_email(doc_path)
    if file_type in EXTRACT_MSG_TYPES:
        return handle_extract_msg(doc_path)
    if file_type in NBFORMAT_TYPES:
        return handle_nbformat(doc_path)
    if file_type in XML_PARSE_TYPES:
        return handle_xml_parse(doc_path)
    if file_type in XLRD_TYPES:
        return handle_xlrd(doc_path)
    if file_type in ANTIWORD_TYPES:
        return handle_antiword(doc_path)
    if file_type in LIBREOFFICE_TYPES:
        return handle_libreoffice(doc_path)
    if file_type in ODS_TYPES:
        return handle_ods(doc_path)
    raise ValueError(f"No handler for file type: {file_type}")


# ---------------------------------------------------------------------------
# Core logic
# ---------------------------------------------------------------------------


def get_repo_root() -> Path:
    """Find the repository root directory."""
    current = Path(__file__).resolve().parent
    while current != current.parent:
        if (current / "Cargo.toml").exists() and (current / "test_documents").exists():
            return current
        current = current.parent
    raise RuntimeError("Could not find repository root")


def collect_fixtures(fixtures_dir: Path) -> list[Path]:
    """Recursively collect all fixture JSON files."""
    return sorted(fixtures_dir.rglob("*.json"))


def load_mapping(repo_root: Path) -> dict[str, str]:
    """Load the existing ground truth mapping."""
    mapping_file = repo_root / "test_documents" / "ground_truth" / "ground_truth_mapping.json"
    if mapping_file.exists():
        with open(mapping_file) as f:
            return json.load(f)
    return {}


def save_mapping(repo_root: Path, mapping: dict[str, str]) -> None:
    """Save the ground truth mapping (sorted keys)."""
    mapping_file = repo_root / "test_documents" / "ground_truth" / "ground_truth_mapping.json"
    sorted_mapping = dict(sorted(mapping.items()))
    with open(mapping_file, "w") as f:
        json.dump(sorted_mapping, f, indent=2)
        f.write("\n")


def make_mapping_key(fixture_path: Path, fixtures_dir: Path) -> str:
    """Generate a unique mapping key from the fixture path.

    For top-level fixtures: stem (e.g., 'commonmark_sample')
    For subdir fixtures: subdir/stem (e.g., 'md/duck.md' from md/duck.md.json)
    """
    rel = fixture_path.relative_to(fixtures_dir)
    parts = rel.parts
    if len(parts) > 1:
        return f"{parts[0]}/{fixture_path.stem}"
    return fixture_path.stem


def process_fixture(
    fixture_path: Path,
    repo_root: Path,
    fixtures_dir: Path,
    mapping: dict[str, str],
    dry_run: bool,
    force: bool,
    stats: dict[str, int],
) -> None:
    """Process a single fixture: generate ground truth, patch fixture, update mapping."""
    with open(fixture_path) as f:
        fixture = json.load(f)

    file_type = fixture.get("file_type", "")

    # Skip excluded types
    if file_type in EXCLUDED_TYPES:
        stats["skipped_excluded"] += 1
        return

    # Skip unhandled types
    if file_type not in ALL_HANDLED_TYPES:
        print(f"  SKIP (unhandled type): {fixture_path.name} ({file_type})")
        stats["skipped_unhandled"] += 1
        return

    # Skip if already has ground truth (unless --force)
    if fixture.get("ground_truth") and not force:
        stats["skipped_existing"] += 1
        return

    # Resolve document path
    doc_rel = fixture.get("document", "")
    if not doc_rel:
        print(f"  SKIP (no document): {fixture_path.name}")
        stats["skipped_no_doc"] += 1
        return

    doc_path = (fixture_path.parent / doc_rel).resolve()
    if not doc_path.exists():
        print(f"  SKIP (doc not found): {fixture_path.name} -> {doc_path}")
        stats["skipped_missing_doc"] += 1
        return

    # Determine ground truth output path
    gt_dir = repo_root / "test_documents" / "ground_truth" / file_type
    gt_filename = fixture_path.stem + ".txt"
    gt_path = gt_dir / gt_filename

    # Compute relative path from fixture to ground truth
    gt_rel = os.path.relpath(gt_path, fixture_path.parent)

    # Mapping key
    mapping_key = make_mapping_key(fixture_path, fixtures_dir)

    if dry_run:
        print(f"  [DRY RUN] {fixture_path.name} ({file_type})")
        print(f"    doc: {doc_path}")
        print(f"    gt:  {gt_path}")
        print(f"    key: {mapping_key}")
        stats["would_generate"] += 1
        return

    # Extract text
    try:
        text = extract_text(doc_path, file_type)
    except Exception as e:
        print(f"  ERROR extracting {fixture_path.name}: {e}")
        stats["errors"] += 1
        return

    # Write ground truth file
    gt_dir.mkdir(parents=True, exist_ok=True)
    gt_path.write_text(text, encoding="utf-8")

    # Patch fixture JSON
    fixture["ground_truth"] = {
        "text_file": gt_rel,
        "source": get_source_type(file_type),
    }
    with open(fixture_path, "w") as f:
        json.dump(fixture, f, indent=2)
        f.write("\n")

    # Update mapping
    gt_mapping_path = str(gt_path.relative_to(repo_root))
    mapping[mapping_key] = gt_mapping_path

    stats["generated"] += 1


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate ground truth for benchmark fixtures")
    parser.add_argument("--dry-run", action="store_true", help="Print planned actions without writing")
    parser.add_argument("--format-filter", type=str, default="", help="Comma-separated file types to process")
    parser.add_argument("--force", action="store_true", help="Regenerate even if ground truth exists")
    parser.add_argument("--skip-types", type=str, default="", help="Comma-separated file types to skip")
    args = parser.parse_args()

    repo_root = get_repo_root()
    fixtures_dir = repo_root / "tools" / "benchmark-harness" / "fixtures"

    print(f"Repository root: {repo_root}")
    print(f"Fixtures dir: {fixtures_dir}")
    if args.dry_run:
        print("DRY RUN MODE - no files will be written\n")

    format_filter = set(args.format_filter.split(",")) if args.format_filter else None
    skip_types = set(args.skip_types.split(",")) if args.skip_types else set()

    # Load existing mapping
    mapping = load_mapping(repo_root)
    initial_mapping_size = len(mapping)

    # Collect and process fixtures
    fixture_paths = collect_fixtures(fixtures_dir)
    print(f"Found {len(fixture_paths)} fixture files\n")

    stats: dict[str, int] = {
        "generated": 0,
        "would_generate": 0,
        "skipped_existing": 0,
        "skipped_excluded": 0,
        "skipped_unhandled": 0,
        "skipped_no_doc": 0,
        "skipped_missing_doc": 0,
        "errors": 0,
    }

    for fixture_path in fixture_paths:
        # Load to check file type for filtering
        try:
            with open(fixture_path) as f:
                fixture_data = json.load(f)
        except (json.JSONDecodeError, OSError) as e:
            print(f"  ERROR reading {fixture_path.name}: {e}")
            stats["errors"] += 1
            continue

        file_type = fixture_data.get("file_type", "")
        if format_filter and file_type not in format_filter:
            continue
        if file_type in skip_types:
            continue

        process_fixture(fixture_path, repo_root, fixtures_dir, mapping, args.dry_run, args.force, stats)

    # Save mapping
    if not args.dry_run and stats["generated"] > 0:
        save_mapping(repo_root, mapping)
        new_entries = len(mapping) - initial_mapping_size
        print(f"\nUpdated ground_truth_mapping.json: {new_entries} new entries (total: {len(mapping)})")

    # Print summary
    print(f"\n{'=' * 50}")
    print("Summary:")
    print(f"  Generated:         {stats['generated']}")
    if args.dry_run:
        print(f"  Would generate:    {stats['would_generate']}")
    print(f"  Skipped (existing): {stats['skipped_existing']}")
    print(f"  Skipped (excluded): {stats['skipped_excluded']}")
    print(f"  Skipped (unhandled): {stats['skipped_unhandled']}")
    print(f"  Skipped (no doc):   {stats['skipped_no_doc']}")
    print(f"  Skipped (missing):  {stats['skipped_missing_doc']}")
    print(f"  Errors:            {stats['errors']}")

    return 1 if stats["errors"] > 0 else 0


if __name__ == "__main__":
    sys.exit(main())
